from pptx import Presentation
from pptx.util import Inches

def export_ppt(nodes, edges):

    prs=Presentation()
    slide=prs.slides.add_slide(prs.slide_layouts[6])

    shapes={}
    x=1
    y=1

    for n,t in nodes:

        box=slide.shapes.add_shape(
            1,
            Inches(x),
            Inches(y),
            Inches(2),
            Inches(1)
        )

        box.text=n
        shapes[n]=box
        x+=2.5

    path="dfd.pptx"
    prs.save(path)

    return path